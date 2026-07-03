from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
slide_layout = prs.slide_layouts[6] # Layout em branco

COR_PRINCIPAL = RGBColor(16, 44, 87)
COR_DESTAQUE = RGBColor(231, 76, 60)
COR_TEXTO = RGBColor(44, 62, 80)

def formatar_titulo(slide, texto, cor=COR_PRINCIPAL, tamanho=40):
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = texto
    p.font.bold = True
    p.font.size = Pt(tamanho)
    p.font.color.rgb = cor
    p.font.name = "Arial"

def adicionar_topicos(slide, lista_topicos, top=2.0):
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(lista_topicos):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COR_TEXTO
        p.font.name = "Arial"
        p.space_after = Pt(12)

# Slide 1: Capa
slide = prs.slides.add_slide(slide_layout)
formatar_titulo(slide, "ChurnGuard AI", COR_PRINCIPAL, 54)
adicionar_topicos(slide, ["Sistema Inteligente de Previsão de Churn Baseado em LTV", "Apresentação Executiva do Projeto", "Desenvolvido por: [Seu Nome]"], top=2.5)

# Slide 2: O Problema
slide = prs.slides.add_slide(slide_layout)
formatar_titulo(slide, "1. A Maior Dor das Empresas Atuais", COR_DESTAQUE)
adicionar_topicos(slide, ["• O Custo de Aquisição de Clientes (CAC) é até 7x maior do que o custo de retenção.", "• Empresas SaaS e E-commerces perdem receita de forma silenciosa sem previsibilidade.", "• Abordagens tradicionais reagem tarde: agem apenas após o cancelamento do contrato.", "• Solução: Antecipar o comportamento de risco para blindar o faturamento mensal (MRR)."])

# Slide 3: Arquitetura Técnica
slide = prs.slides.add_slide(slide_layout)
formatar_titulo(slide, "2. Arquitetura de Produção Multidisciplinar")
adicionar_topicos(slide, ["• Engenharia de Dados: Ingestão automatizada separada em camadas (Raw, Processed, Gold).", "• Ciência de Dados: Algoritmo Random Forest com 94.00% de acurácia de predição.", "• Desenvolvimento: API assíncrona robusta em FastAPI e validação com Pydantic.", "• Análise de Dados: Dashboard interativo em Streamlit focado em KPIs de ROI.", "• Infraestrutura: Toda a aplicação empacotada em Docker Containers."])

# Slide 4: Resultados Práticos
slide = prs.slides.add_slide(slide_layout)
formatar_titulo(slide, "3. Impacto e Métricas de Negócio")
adicionar_topicos(slide, ["• Visibilidade Financeira: Tradução de métricas técnicas em MRR Crítico em Risco.", "• Ação Imediata: Geração de listas automatizadas de prioridade para Customer Success.", "• IA Explicável: Modelagem baseada na probabilidade exata de evasão de cada cliente.", "• Decisão Estratégica: Painel unificado para tomadas de decisão de Diretores e CFOs."])

prs.save("docs/apresentacao_churn_guard_ai.pptx")
print("✅ PowerPoint gerado com sucesso em: docs/apresentacao_churn_guard_ai.pptx")

